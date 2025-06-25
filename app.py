import streamlit as st
import google.generativeai as genai
import fitz
import traceback
import time
import re
from io import BytesIO
from typing import Dict, Any, Optional, List
import json
from dataclasses import dataclass

# 需要安装的依赖包
# pip install streamlit PyMuPDF python-pptx google-generativeai

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    import colorsys
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# --- 提示词模板 (保持不变) ---

# ## 大纲生成器 (保持不变) ##
OUTLINE_GENERATION_PROMPT_TEMPLATE = """
角色 (Role):
你是一位顶级的学术汇报设计师和内容策略师，同时具备出色的**"无图化设计" (Graphic-less Design)** 思维。你精通将复杂的学术论文转化为结构化、视觉化的演示文稿（PPT），并且擅长使用CSS样式、布局和文本符号来创造清晰、优雅的视觉效果，以最大限度地减少对外部图片或复杂SVG的依赖。

核心任务 (Core Task):
分析用户上传的学术文档，并生成一个结构化的、逐页的演示文稿大纲。这个大纲需要直接映射到一个预设的HTML汇报模板中。你需要决定将文档内容划分成多少页PPT是合适的（通常在10-15页之间），并为每一页设计内容和可通过简单代码实现的视觉元素。

关键原则 (Guiding Principles):
一页一核心: 每张幻灯片只传达一个核心观点。
化繁为简: 将长句转化为精炼的要点、短语或关键词。
逻辑流畅: 遵循标准的学术汇报逻辑（引言 -> 方法 -> 结果 -> 结论）。
提取关键: 直接从原文中提取关键术语、数据、结果和引用。
CSS优先，无图为主 (CSS First, Image-Free): 这是最重要的原则。 你必须优先考虑那些可以通过纯CSS、Emoji或基本HTML结构实现的视觉元素。绝对禁止建议使用外部图片链接。 仅在极其必要且形状极其简单（如圆形、箭头）的情况下，才建议使用SVG路径。

输出格式 (Required Output Format):
你必须严格按照以下Markdown格式为每一页幻灯片生成内容。使用 --- 分隔每一页。

---
**Slide:** [幻灯片页码，从1开始]
**Title:** [幻灯片标题]
**Purpose:** [幻灯片目的/类型，从以下选项中选择: Title, Overview, Background, Methodology, Data, Results, Analysis, Discussion, Conclusion, Future_Work, Acknowledgements]
**Content:**
- [要点1：简洁的短句或短语]
- [要点2：**加粗**关键术语]
- [要点3：直接引用原文的一句核心观点]
**Visual:**
  - **Type:** [从以下视觉类型中选择: `Symbol`, `Process`, `Chart`, `Table`, `Quote`, `Comparison`, `List`, `Text_Only`]
  - **Data:** [根据选择的Type提供结构化数据。这是最关键的部分，格式见下方说明。]
---

Visual.Data 格式说明 (已优化):
Type: Symbol
Data:
symbol: [一个Unicode Emoji表情符号]
text: [符号旁边的简短说明文字]
color_hint: [一个CSS颜色提示]

Type: Process
Data:
steps: [一个JSON数组]
style: [numbered-list, chevron-arrow]

Type: Chart
Data:
chart_type: [bar, line, pie]
title: [图表标题]
data_summary: [对图表核心数据的文字描述]

Type: Table
Data:
caption: [表格标题]
headers: [一个JSON数组]
rows: [一个包含多行数据的JSON数组]

Type: Quote
Data:
text: [引用的核心文本]
source: [引用来源]

Type: Comparison
Data:
item1_title: [对比项1的标题]
item1_points: [一个JSON数组]
item2_title: [对比项2的标题]
item2_points: [一个JSON数组]

Type: List 或 Type: Text_Only
Data: null

指令 (Instruction):
现在，请分析用户上传的这份学术文档。严格遵循以上所有规则和**"无图化设计"原则，为其生成一份完整的、逻辑清晰的、强调使用简单符号和CSS**进行视觉呈现的学术演示文稿大纲。请开始。
"""

# ## 代码融合器 (终极强化版 - 结合原始流程精华) ##
CODE_GENERATION_PROMPT_TEMPLATE = """
角色 (Role):
你是一位精通HTML、CSS和JavaScript的前端开发专家，拥有像素级的代码保真能力。你的核心任务是将结构化的Markdown大纲，无损地、精确地与一个预定义的HTML模板相结合，动态生成最终的、可直接运行的、高度专业的HTML文件。你对细节有极高的要求，尤其是在处理图像资源和数据可视化占位方面。

核心任务 (Core Task):
你将收到两份输入：
1. **PPT大纲 (PPT Outline):** 一份由AI预先生成的、结构化的Markdown文件。
2. **HTML模板 (HTML Template):** 一个完整的HTML文件，包含了所有必须的样式、脚本和关键资源（如Base64编码的校徽）。

你的任务是：
1. **解析大纲:** 逐页解析PPT大纲中的所有字段（Slide、Title、Purpose、Content、Visual等）。
2. **动态生成幻灯片:** 根据解析出的数据，为每一页幻灯片生成对应的HTML `<section>` 元素，并应用正确的CSS类。
3. **智能渲染视觉元素:**
   - **对于图表 (Visual.Type: Chart):** 绝不在页面上显示"占位符"字样。你应该在图表区域内，使用优雅的排版，将大纲中提供的 `Visual.Data.data_summary` (数据摘要文字) 直接展示出来。这为演讲者提供了一个讨论数据的起点，而不是一个空洞的占位符。
   - **对于符号 (Visual.Type: Symbol):** 将大纲中指定的Emoji符号 (`Visual.Data.symbol`) 直接作为文本插入到HTML中，并可选择性地使用 `Visual.Data.color_hint` 作为内联样式的颜色。
   - **对于其他类型:** 根据Visual.Type和Visual.Data智能生成相应的HTML结构。
4. **【最高优先级 - 铁律】保护关键资源:** 在整合代码时，必须完整、无误地保留HTML模板中所有的：
   - 整个`<head>`标签，包含所有的`<link>`和`<style>`
   - 整个`<script>`标签及其内部所有的JavaScript代码
   - 所有导航控件、页码指示器等非幻灯片内容
   - **特别重要:** 所有`<img>`标签及其`src`属性，尤其是那些包含 `data:image/svg+xml;base64,...` 的长字符串。绝不允许对这些资源链接进行任何形式的修改、缩短或删除。
5. **无缝整合:** 确保动态生成的幻灯片数量与底部的缩略图导航和演讲者备注的条目数量完全一致。
6. **代码整洁:** 生成的HTML代码必须有良好的缩进和可读性。

**【绝对禁止 - 输出要求】:**
- 你的最终输出 **绝对不能** 包含任何解释性文字或Markdown代码块标记
- 不要使用```html或```等Markdown代码块标记
- 不要在HTML前后添加任何解释性内容
- 输出必须是一个纯粹的HTML文本，直接以 `<!DOCTYPE html>` 开头，并以 `</html>` 结尾

指令 (Instruction):
以下是用户提供的 **PPT大纲 (PPT Outline)** 和 **HTML模板 (HTML Template)**。请你立即开始工作，严格遵循以上所有规则，特别是保护校徽等关键资源和优雅处理图表占位的指令，将大纲内容与模板代码结合，生成最终的、完整的、专业级的HTML文件。不要提供任何解释或评论，直接输出完整的HTML代码。
"""

# --- 复旦大学模板样式提取器 ---
class FudanTemplateExtractor:
    """专门针对复旦大学HTML模板的精确样式提取器"""
    
    def __init__(self, html_content: str):
        self.html_content = html_content
        self.fudan_colors = {}
        self.fudan_fonts = {}
        self.fudan_spacing = {}
        self._extract_fudan_design_system()
    
    def _extract_fudan_design_system(self):
        """提取复旦模板的完整设计系统"""
        try:
            # 提取CSS变量（复旦设计系统的核心）
            css_vars_pattern = r':root\s*{([^}]+)}'
            css_vars_match = re.search(css_vars_pattern, self.html_content, re.DOTALL)
            
            if css_vars_match:
                css_vars = css_vars_match.group(1)
                
                # 提取复旦色彩系统
                color_patterns = {
                    'fudan_blue': r'--fudan-blue:\s*([^;]+)',
                    'fudan_deep_blue': r'--fudan-deep-blue:\s*([^;]+)', 
                    'fudan_light_blue': r'--fudan-light-blue:\s*([^;]+)',
                    'fudan_white': r'--fudan-white:\s*([^;]+)',
                    'fudan_black': r'--fudan-black:\s*([^;]+)',
                    'fudan_near_black': r'--fudan-near-black:\s*([^;]+)',
                    'fudan_bg_light': r'--fudan-bg-light:\s*([^;]+)',
                    'fudan_accent_gray': r'--fudan-accent-gray:\s*([^;]+)',
                    'fudan_light_gray': r'--fudan-light-gray:\s*([^;]+)'
                }
                
                for key, pattern in color_patterns.items():
                    match = re.search(pattern, css_vars)
                    if match:
                        self.fudan_colors[key] = match.group(1).strip()
                
                # 提取字体系统
                font_patterns = {
                    'serif': r'--font-serif:\s*([^;]+)',
                    'sans': r'--font-sans:\s*([^;]+)'
                }
                
                for key, pattern in font_patterns.items():
                    match = re.search(pattern, css_vars)
                    if match:
                        self.fudan_fonts[key] = match.group(1).strip().replace("'", "").split(',')[0]
            
            # 设置默认复旦样式（如果提取失败）
            if not self.fudan_colors:
                self.fudan_colors = {
                    'fudan_blue': '#0055A2',
                    'fudan_deep_blue': '#003366',
                    'fudan_light_blue': '#A8D8F8',
                    'fudan_white': '#FFFFFF',
                    'fudan_near_black': '#2D3748',
                    'fudan_bg_light': '#F7FAFC',
                    'fudan_light_gray': '#E2E8F0'
                }
            
            if not self.fudan_fonts:
                self.fudan_fonts = {
                    'serif': 'Noto Serif SC',
                    'sans': 'Noto Sans SC'
                }
                
        except Exception as e:
            # 确保有默认的复旦样式
            self.fudan_colors = {
                'fudan_blue': '#0055A2',
                'fudan_deep_blue': '#003366', 
                'fudan_white': '#FFFFFF',
                'fudan_near_black': '#2D3748',
                'fudan_bg_light': '#F7FAFC'
            }
            self.fudan_fonts = {'serif': 'Noto Serif SC', 'sans': 'Noto Sans SC'}

# --- 复旦风格PPT生成器 ---
class FudanStylePPTGenerator:
    """完全按照复旦大学HTML模板风格生成PPT"""
    
    def __init__(self, html_template: str = None):
        self.presentation = None
        self.fudan_extractor = FudanTemplateExtractor(html_template) if html_template else None
    
    def create_presentation(self, outline_data: str) -> BytesIO:
        """根据大纲创建复旦风格PPT"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx库未安装，无法生成PPT文件")
        
        # 创建演示文稿
        self.presentation = Presentation()
        
        # 设置演示文稿的默认样式为复旦风格
        self._setup_fudan_master_styles()
        
        # 解析大纲
        slides_data = self._parse_outline(outline_data)
        
        # 生成复旦风格幻灯片
        for slide_data in slides_data:
            self._create_fudan_style_slide(slide_data)
        
        # 保存到BytesIO
        ppt_buffer = BytesIO()
        self.presentation.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
    
    def _setup_fudan_master_styles(self):
        """设置复旦大学风格的主样式"""
        pass  # PPT母版样式设置比较复杂，这里先跳过
    
    def _parse_outline(self, outline_text: str) -> List[Dict]:
        """解析大纲文本，保持与原版相同"""
        slides = []
        current_slide = {}
        current_content = []
        in_content_section = False
        in_visual_section = False
        
        lines = outline_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line == '---':
                if current_slide:
                    current_slide['content'] = current_content
                    slides.append(current_slide)
                    current_slide = {}
                    current_content = []
                    in_content_section = False
                    in_visual_section = False
                continue
            
            if line.startswith('**Slide:**'):
                current_slide['slide_num'] = line.split('**Slide:**')[1].strip()
            elif line.startswith('**Title:**'):
                current_slide['title'] = line.split('**Title:**')[1].strip()
            elif line.startswith('**Purpose:**'):
                current_slide['purpose'] = line.split('**Purpose:**')[1].strip()
            elif line.startswith('**Content:**'):
                in_content_section = True
                in_visual_section = False
            elif line.startswith('**Visual:**'):
                in_visual_section = True
                in_content_section = False
                current_slide['visual'] = {'type': '', 'data': ''}
            elif line.startswith('  - **Type:**') and in_visual_section:
                current_slide['visual']['type'] = line.split('**Type:**')[1].strip()
            elif line.startswith('  - **Data:**') and in_visual_section:
                current_slide['visual']['data'] = line.split('**Data:**')[1].strip()
            elif line.startswith('- ') and in_content_section:
                # 清理markdown格式
                content_line = line[2:].strip()
                content_line = re.sub(r'\*\*(.*?)\*\*', r'\1', content_line)  # 移除加粗标记
                current_content.append(content_line)
        
        # 添加最后一个幻灯片
        if current_slide:
            current_slide['content'] = current_content
            slides.append(current_slide)
            
        return slides
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """将十六进制颜色转换为RGB"""
        hex_color = hex_color.strip('#')
        if len(hex_color) == 6:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return (0, 85, 162)  # 默认复旦蓝
    
    def _create_fudan_style_slide(self, slide_data: Dict):
        """创建复旦风格的幻灯片"""
        purpose = slide_data.get('purpose', 'Content')
        
        # 根据purpose选择布局
        if purpose == 'Title':
            layout = self.presentation.slide_layouts[0]  # 标题布局
        else:
            layout = self.presentation.slide_layouts[1]  # 内容布局
            
        slide = self.presentation.slides.add_slide(layout)
        
        # 设置复旦风格背景
        self._apply_fudan_background(slide)
        
        # 处理标题
        title_text = slide_data.get('title', '')
        if title_text and slide.shapes.title:
            self._create_fudan_title(slide, title_text, purpose)
        
        # 处理内容
        content = slide_data.get('content', [])
        if content:
            self._create_fudan_content_card(slide, content, title_text)
        
        # 处理视觉元素
        visual = slide_data.get('visual', {})
        if visual and visual.get('type'):
            self._add_fudan_visual_element(slide, visual)
    
    def _apply_fudan_background(self, slide):
        """应用复旦风格背景"""
        try:
            if self.fudan_extractor:
                # 设置浅色背景
                bg_color = self.fudan_extractor.fudan_colors.get('fudan_bg_light', '#F7FAFC')
                background = slide.background
                fill = background.fill
                fill.solid()
                rgb = self._hex_to_rgb(bg_color)
                fill.fore_color.rgb = RGBColor(*rgb)
        except Exception:
            pass
    
    def _create_fudan_title(self, slide, title_text: str, purpose: str):
        """创建复旦风格标题"""
        try:
            title_shape = slide.shapes.title
            title_shape.text = title_text
            
            # 应用复旦风格
            text_frame = title_shape.text_frame
            
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                
                for run in paragraph.runs:
                    # 设置复旦蓝色
                    if self.fudan_extractor:
                        fudan_blue = self.fudan_extractor.fudan_colors.get('fudan_blue', '#0055A2')
                        rgb = self._hex_to_rgb(fudan_blue)
                        run.font.color.rgb = RGBColor(*rgb)
                    
                    # 设置字体
                    if purpose == 'Title':
                        run.font.size = Pt(44)  # 大标题
                        run.font.name = 'Microsoft YaHei'  # Windows中文字体
                    else:
                        run.font.size = Pt(36)  # 普通标题
                        run.font.name = 'Microsoft YaHei'
                    
                    run.font.bold = True
        except Exception:
            pass
    
    def _create_fudan_content_card(self, slide, content: List[str], title_text: str):
        """创建复旦风格的内容卡片"""
        try:
            # 删除默认的内容占位符
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                slide.shapes._spTree.remove(content_placeholder._element)
            
            # 创建复旦风格的内容卡片
            # 计算位置和大小
            left = Inches(1)
            top = Inches(2.5)  # 在标题下方
            width = Inches(8.5)
            height = Inches(5)
            
            # 创建背景卡片（模拟HTML中的research-card）
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            
            # 设置卡片样式（模拟复旦HTML样式）
            if self.fudan_extractor:
                # 白色背景
                card_fill = card_shape.fill
                card_fill.solid()
                white_color = self.fudan_extractor.fudan_colors.get('fudan_white', '#FFFFFF')
                rgb = self._hex_to_rgb(white_color)
                card_fill.fore_color.rgb = RGBColor(*rgb)
                
                # 设置边框（模拟顶部蓝色边框）
                line = card_shape.line
                line.color.rgb = RGBColor(*self._hex_to_rgb(
                    self.fudan_extractor.fudan_colors.get('fudan_blue', '#0055A2')
                ))
                line.width = Pt(4)
            
            # 创建内容文本框
            text_left = left + Inches(0.5)
            text_top = top + Inches(0.5)
            text_width = width - Inches(1)
            text_height = height - Inches(1)
            
            textbox = slide.shapes.add_textbox(
                text_left, text_top, text_width, text_height
            )
            text_frame = textbox.text_frame
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            
            # 添加内容
            if content:
                # 第一段作为主文本
                text_frame.text = content[0]
                
                # 设置第一段样式
                first_paragraph = text_frame.paragraphs[0]
                self._apply_fudan_text_style(first_paragraph, is_main_text=True)
                
                # 添加其他段落
                for item in content[1:]:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"  # 添加项目符号
                    self._apply_fudan_text_style(p, is_bullet=True)
                    
        except Exception as e:
            # 如果创建卡片失败，回退到简单文本
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                
                if content:
                    tf.text = content[0]
                    for item in content[1:]:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
    
    def _apply_fudan_text_style(self, paragraph, is_main_text=False, is_bullet=False):
        """应用复旦文本样式"""
        try:
            for run in paragraph.runs:
                # 设置字体
                run.font.name = 'Microsoft YaHei'
                
                if is_main_text:
                    run.font.size = Pt(20)
                    # 深色文字
                    if self.fudan_extractor:
                        text_color = self.fudan_extractor.fudan_colors.get('fudan_near_black', '#2D3748')
                        rgb = self._hex_to_rgb(text_color)
                        run.font.color.rgb = RGBColor(*rgb)
                elif is_bullet:
                    run.font.size = Pt(18)
                    # 稍浅的文字
                    if self.fudan_extractor:
                        text_color = self.fudan_extractor.fudan_colors.get('fudan_near_black', '#2D3748')
                        rgb = self._hex_to_rgb(text_color)
                        run.font.color.rgb = RGBColor(*rgb)
                else:
                    run.font.size = Pt(18)
                
            # 设置段落间距
            paragraph.space_after = Pt(12)
            
        except Exception:
            pass
    
    def _add_fudan_visual_element(self, slide, visual):
        """添加复旦风格的视觉元素"""
        visual_type = visual.get('type', '').strip('`')
        
        if visual_type == 'Symbol':
            self._add_fudan_symbol_element(slide, visual)
        elif visual_type == 'Chart':
            self._add_fudan_chart_element(slide, visual)
        elif visual_type == 'Table':
            self._add_fudan_table_element(slide, visual)
    
    def _add_fudan_symbol_element(self, slide, visual):
        """添加复旦风格符号元素"""
        try:
            data_str = visual.get('data', '')
            if 'symbol:' in data_str:
                # 解析symbol
                symbol_match = re.search(r'symbol:\s*([^\n]+)', data_str)
                if symbol_match:
                    symbol = symbol_match.group(1).strip()
                    
                    # 在右上角添加符号（复旦风格位置）
                    left = Inches(8.5)
                    top = Inches(1)
                    width = Inches(1)
                    height = Inches(1)
                    
                    # 创建符号背景圆圈（模拟复旦seal-style）
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        left, top, width, height
                    )
                    
                    # 设置复旦蓝色背景
                    if self.fudan_extractor:
                        fill = circle.fill
                        fill.solid()
                        fudan_blue = self.fudan_extractor.fudan_colors.get('fudan_blue', '#0055A2')
                        rgb = self._hex_to_rgb(fudan_blue)
                        fill.fore_color.rgb = RGBColor(*rgb)
                    
                    # 添加符号文本
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = symbol
                    
                    # 设置符号样式
                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = RGBColor(255, 255, 255)  # 白色
        except Exception:
            pass
    
    def _add_fudan_chart_element(self, slide, visual):
        """添加复旦风格图表元素"""
        try:
            data_str = visual.get('data', '')
            if 'data_summary:' in data_str:
                # 提取数据摘要
                summary_match = re.search(r'data_summary:\s*([^\n]+)', data_str)
                if summary_match:
                    summary = summary_match.group(1).strip()
                    
                    # 创建复旦风格的数据卡片
                    left = Inches(1)
                    top = Inches(6)
                    width = Inches(8.5)
                    height = Inches(1.5)
                    
                    # 创建背景卡片
                    card_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        left, top, width, height
                    )
                    
                    # 设置卡片样式
                    if self.fudan_extractor:
                        # 浅蓝色背景
                        card_fill = card_shape.fill
                        card_fill.solid()
                        light_blue = self.fudan_extractor.fudan_colors.get('fudan_light_blue', '#A8D8F8')
                        rgb = self._hex_to_rgb(light_blue)
                        card_fill.fore_color.rgb = RGBColor(*rgb)
                    
                    # 添加数据摘要文本
                    textbox = slide.shapes.add_textbox(
                        left + Inches(0.2), top + Inches(0.2), 
                        width - Inches(0.4), height - Inches(0.4)
                    )
                    text_frame = textbox.text_frame
                    text_frame.text = f"📊 数据要点: {summary}"
                    
                    # 设置文本样式
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Microsoft YaHei'
                            run.font.size = Pt(16)
                            run.font.bold = True
                            if self.fudan_extractor:
                                fudan_blue = self.fudan_extractor.fudan_colors.get('fudan_deep_blue', '#003366')
                                rgb = self._hex_to_rgb(fudan_blue)
                                run.font.color.rgb = RGBColor(*rgb)
        except Exception:
            pass
    
    def _add_fudan_table_element(self, slide, visual):
        """添加复旦风格表格元素"""
        try:
            data_str = visual.get('data', '') 
            if 'caption:' in data_str:
                caption_match = re.search(r'caption:\s*([^\n]+)', data_str)
                if caption_match:
                    caption = caption_match.group(1).strip()
                    
                    # 添加表格标题卡片
                    left = Inches(1)
                    top = Inches(5.5)
                    width = Inches(8.5)
                    height = Inches(0.8)
                    
                    # 创建标题背景
                    card_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        left, top, width, height
                    )
                    
                    # 设置复旦风格
                    if self.fudan_extractor:
                        card_fill = card_shape.fill
                        card_fill.solid()
                        fudan_gray = self.fudan_extractor.fudan_colors.get('fudan_light_gray', '#E2E8F0')
                        rgb = self._hex_to_rgb(fudan_gray)
                        card_fill.fore_color.rgb = RGBColor(*rgb)
                    
                    # 添加标题文本
                    textbox = slide.shapes.add_textbox(
                        left + Inches(0.2), top + Inches(0.1),
                        width - Inches(0.4), height - Inches(0.2)
                    )
                    text_frame = textbox.text_frame
                    text_frame.text = f"📋 {caption}"
                    
                    # 设置文本样式
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Microsoft YaHei'
                            run.font.size = Pt(18)
                            run.font.bold = True
                            if self.fudan_extractor:
                                fudan_blue = self.fudan_extractor.fudan_colors.get('fudan_blue', '#0055A2')
                                rgb = self._hex_to_rgb(fudan_blue)
                                run.font.color.rgb = RGBColor(*rgb)
        except Exception:
            pass

# --- 所有Agent函数 (保持与原始版本完全一致) ---
def parse_pdf(uploaded_file, debug_log_container):
    try:
        file_bytes = uploaded_file.getvalue()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        full_text = "".join(page.get_text() + "\n" for page in doc)
        debug_log_container.write(f"✅ PDF解析成功。总计 {len(full_text):,} 个字符。")
        return full_text
    except Exception:
        debug_log_container.error(f"PDF解析时出现异常: {traceback.format_exc()}")
        return None

def validate_model(api_key, model_name, debug_log_container):
    try:
        if not model_name or not model_name.strip(): return False
        genai.configure(api_key=api_key)
        available_models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        if f"models/{model_name}" in available_models:
            debug_log_container.success(f"✅ 模型 `{model_name}` 验证通过！")
            return True
        else:
            st.error(f"**模型验证失败!** `{model_name}` 不存在。")
            return False
    except Exception:
        st.error(f"**API Key验证失败!**")
        debug_log_container.error(f"验证API Key时异常: {traceback.format_exc()}")
        return False

def call_gemini(api_key, prompt_text, ui_placeholder, model_name, debug_log_container):
    try:
        debug_log_container.write(f"--- \n准备调用AI: `{model_name}`...")
        debug_log_container.write(f"**发送的Prompt长度:** `{len(prompt_text):,}` 字符")
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(model_name)
        
        collected_chunks = []
        def stream_and_collect(stream):
            for chunk in stream:
                if hasattr(chunk, 'text'):
                    text_part = chunk.text
                    collected_chunks.append(text_part)
                    yield text_part

        # 只有在提供了UI占位符时才进行流式写入
        if ui_placeholder:
            response_stream = model.generate_content(prompt_text, stream=True)
            ui_placeholder.write_stream(stream_and_collect(response_stream))
        else:
            # 如果不提供UI占位符，则直接生成，避免在UI上产生不必要的输出
            response = model.generate_content(prompt_text)
            if hasattr(response, 'text'):
                collected_chunks.append(response.text)
        
        full_response_str = "".join(collected_chunks)
        debug_log_container.write(f"✅ AI响应成功完成。收集到 {len(full_response_str):,} 个字符。")
        return full_response_str
    except Exception as e:
        error_message = f"🚨 **AI调用失败!** 请检查调试日志。\n\n**错误详情:** {e}"
        if ui_placeholder:
            ui_placeholder.error(error_message)
        else:
            st.error(error_message)
        debug_log_container.error(f"--- AI调用时发生严重错误 ---\n{traceback.format_exc()}")
        return None

def extract_clean_outline(raw_output, debug_log_container):
    try:
        match = re.search(r"\*\*\s*Slide\s*:\s*\*\*", raw_output)
        if not match:
            debug_log_container.error("❌ 在AI响应中未能找到任何`**Slide:**`标记。")
            return None
        first_slide_pos = match.start()
        last_divider_pos = raw_output.rfind("---", 0, first_slide_pos)
        cleaned_outline = raw_output[last_divider_pos:] if last_divider_pos != -1 else raw_output[first_slide_pos:]
        cleaned_outline = cleaned_outline.strip()
        if cleaned_outline.count("**Title:**") < 3:
            debug_log_container.warning("⚠️ 提取出的大纲结构不完整。")
        debug_log_container.success(f"✅ 已智能识别并提取出大纲内容。")
        return cleaned_outline
    except Exception:
        debug_log_container.error(f"提取大纲时发生意外错误: {traceback.format_exc()}")
        return None

# ## 强化版最终清理函数 - 修复HTML结束标签问题 ##
def final_cleanup(raw_html, debug_log_container):
    """
    对最终的HTML进行强力清理，彻底解决HTML结束标签问题。
    """
    try:
        debug_log_container.write(f"开始清理HTML，原始长度: {len(raw_html):,} 字符")
        
        # 1. 先清理明显的markdown标记
        cleaned_html = raw_html
        markdown_patterns = [
            r'```html\s*',
            r'```\s*$',
            r'^```.*?\n',
        ]
        
        for pattern in markdown_patterns:
            cleaned_html = re.sub(pattern, '', cleaned_html, flags=re.MULTILINE)
        
        # 2. 寻找HTML文档的真正起点
        html_start_pos = cleaned_html.find("<!DOCTYPE html>")
        if html_start_pos == -1:
            debug_log_container.warning("⚠️ 未找到`<!DOCTYPE html>`，尝试寻找`<html`标签")
            html_start_pos = cleaned_html.find("<html")
            if html_start_pos == -1:
                debug_log_container.error("❌ 未找到HTML起始标签")
                return None
        
        # 3. 寻找HTML文档的结束位置 - 改进算法
        html_end_pos = -1
        
        # 先尝试找到最后一个</html>
        html_end_matches = list(re.finditer(r'</html>', cleaned_html, re.IGNORECASE))
        if html_end_matches:
            html_end_pos = html_end_matches[-1].end()
            debug_log_container.write(f"找到HTML结束标签，位置: {html_end_pos}")
        else:
            # 如果没找到</html>，检查是否有</body>，然后手动添加</html>
            body_end_matches = list(re.finditer(r'</body>', cleaned_html, re.IGNORECASE))
            if body_end_matches:
                body_end_pos = body_end_matches[-1].end()
                # 在</body>后添加</html>
                cleaned_html = cleaned_html[:body_end_pos] + "\n</html>" + cleaned_html[body_end_pos:]
                html_end_pos = body_end_pos + 8  # 8 = len("\n</html>")
                debug_log_container.write(f"未找到</html>，在</body>后添加，新位置: {html_end_pos}")
            else:
                # 最后的fallback：在最后添加</html>
                cleaned_html += "\n</html>"
                html_end_pos = len(cleaned_html)
                debug_log_container.write(f"未找到</body>，在末尾添加</html>，位置: {html_end_pos}")
        
        # 4. 提取HTML内容
        html_content = cleaned_html[html_start_pos:html_end_pos].strip()
        
        # 5. 最终验证和修复
        if not html_content.endswith("</html>"):
            html_content += "\n</html>"
            debug_log_container.write("添加缺失的</html>标签")
        
        # 6. 基本格式验证
        if not (html_content.startswith("<!DOCTYPE html>") or html_content.startswith("<html")):
            debug_log_container.error("❌ 清理后的HTML格式不正确")
            return None
        
        debug_log_container.success(f"✅ HTML清理完成！最终长度: {len(html_content):,} 字符")
        debug_log_container.write(f"HTML开头: {html_content[:100]}...")
        debug_log_container.write(f"HTML结尾: ...{html_content[-100:]}")
        
        return html_content
        
    except Exception as e:
        debug_log_container.error(f"最终清理时出错: {traceback.format_exc()}")
        return None

# --- 配置区域 (用户可预设默认API Key) ---
# 🔑 在下方引号内填入您的Gemini API Key，避免每次手动输入
DEFAULT_GEMINI_API_KEY = ""  # 在这里填入您的API Key

# --- Streamlit UI ---
st.set_page_config(page_title="AI学术汇报生成器", page_icon="🎓", layout="wide")
st.title("🎓 AI学术汇报一键生成器 (增强版)")

with st.sidebar:
    st.header("⚙️ 配置")
    # 如果有默认API Key，则预填充，否则为空
    default_key = DEFAULT_GEMINI_API_KEY if DEFAULT_GEMINI_API_KEY.strip() else ""
    api_key = st.text_input("请输入您的Google Gemini API Key", 
                           value=default_key, 
                           type="password",
                           help="💡 提示：您可以在代码顶部的 DEFAULT_GEMINI_API_KEY 中预设API Key")
    model_options =  ['gemini-2.5-flash', 'gemini-2.0-flash', 'gemini-1.5-flash']
    selected_model = st.selectbox("选择AI模型", model_options, index=0)
    
    st.divider()
    st.header("📄 输出格式")
    output_formats = st.multiselect(
        "选择输出格式",
        ["HTML演示文稿", "PPT文件"],
        default=["HTML演示文稿"]
    )

col1, col2 = st.columns(2)
with col1: 
    pdf_file = st.file_uploader("1. 上传您的学术论文 (.pdf)", type=['pdf'])
with col2: 
    html_template = st.file_uploader("2. 上传您的HTML模板", type=['html'])

# 存储生成结果
if 'results' not in st.session_state: 
    st.session_state.results = {}

# --- 主流程 ---
if st.button("🚀 开始生成汇报", use_container_width=True, 
             disabled=(not api_key or not pdf_file or not output_formats or 
                      ("HTML演示文稿" in output_formats and not html_template))):
    
    st.session_state.results = {}
    progress_container = st.container()
    progress_text = progress_container.empty()
    progress_bar = progress_container.progress(0)
    
    with st.expander("🐞 **调试日志 (点击展开查看详细流程)**", expanded=False):
        debug_log_container = st.container()

    if not validate_model(api_key, selected_model, debug_log_container): 
        st.stop()
    progress_bar.progress(5)

    paper_text = parse_pdf(pdf_file, debug_log_container)
    if paper_text:
        progress_bar.progress(10)
        
        progress_text.text(f"步骤 1/3: 正在深度分析生成大纲...")
        prompt_for_outline = OUTLINE_GENERATION_PROMPT_TEMPLATE + "\n\n--- 学术文档全文 ---\n" + paper_text
        outline_placeholder = st.empty()
        markdown_outline = call_gemini(api_key, prompt_for_outline, outline_placeholder, selected_model, debug_log_container)
        
        if markdown_outline:
            progress_bar.progress(40)
            outline_placeholder.empty()

            progress_text.text(f"步骤 2/3: 正在智能识别并清洗大纲...")
            cleaned_outline = extract_clean_outline(markdown_outline, debug_log_container)

            if cleaned_outline:
                progress_bar.progress(50)
                
                # 生成HTML
                if "HTML演示文稿" in output_formats:
                    progress_text.text(f"步骤 3a/3: 正在融合大纲与模板生成HTML文件...")
                    st.info("ℹ️ AI正在执行最终的全文重写，这可能需要一些时间...")
                    
                    template_code = html_template.getvalue().decode("utf-8")
                    
                    final_prompt = "".join([
                        CODE_GENERATION_PROMPT_TEMPLATE, 
                        "\n\n--- PPT Outline ---\n", 
                        cleaned_outline, 
                        "\n\n--- HTML Template ---\n", 
                        template_code
                    ])
                    
                    with st.spinner("AI正在生成最终HTML，请稍候..."):
                        final_html_raw = call_gemini(api_key, final_prompt, None, selected_model, debug_log_container)

                    if final_html_raw:
                        final_html_code = final_cleanup(final_html_raw, debug_log_container)

                        if final_html_code and "</html>" in final_html_code.lower():
                            debug_log_container.success(f"✅ 最终HTML生成并清理成功！")
                            st.session_state.results['html'] = final_html_code
                        else:
                            st.error("AI未能生成有效的最终HTML文件。请检查调试日志。")
                    else:
                        st.error("AI未能生成最终HTML内容。")
                
                progress_bar.progress(70)
                
                # 生成PPT
                if "PPT文件" in output_formats:
                    if not PPTX_AVAILABLE:
                        st.error("PPT生成需要安装python-pptx库: pip install python-pptx")
                    else:
                        progress_text.text(f"步骤 3b/3: 正在生成PPT文件...")
                        
                        try:
                            # 传入HTML模板用于样式提取
                            template_code = html_template.getvalue().decode("utf-8") if html_template else None
                            ppt_generator = FudanStylePPTGenerator(template_code)
                            ppt_buffer = ppt_generator.create_presentation(cleaned_outline)
                            st.session_state.results['ppt'] = ppt_buffer.getvalue()
                            debug_log_container.success("✅ 复旦风格PPT生成成功！")
                        except Exception as e:
                            st.error(f"PPT生成失败: {e}")
                            debug_log_container.error(f"PPT生成错误: {traceback.format_exc()}")
                
                progress_bar.progress(100)
                progress_text.text(f"🎉 全部完成！")
            else:
                st.error("无法从AI响应中提取出有效的大纲。")

# 下载区域
if st.session_state.results:
    st.header("📥 下载生成的文件")
    
    col1, col2 = st.columns(2)
    
    with col1:
        if 'html' in st.session_state.results:
            st.download_button(
                label="📥 下载HTML演示文稿",
                data=st.session_state.results['html'].encode('utf-8'),
                file_name='presentation.html',
                mime='text/html',
                use_container_width=True
            )
    
    with col2:
        if 'ppt' in st.session_state.results:
            st.download_button(
                label="📥 下载PPT文件",
                data=st.session_state.results['ppt'],
                file_name='presentation.pptx',
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                use_container_width=True
            )

# 依赖说明
st.sidebar.divider()
st.sidebar.header("📦 依赖库状态")
st.sidebar.write("✅ PyMuPDF")
st.sidebar.write("✅ google-generativeai")
st.sidebar.write("✅ python-pptx" if PPTX_AVAILABLE else "❌ python-pptx")

if not PPTX_AVAILABLE:
    st.sidebar.info("💡 安装python-pptx以启用PPT生成功能:\npip install python-pptx")